#!/usr/bin/env python3
"""Builds MoonDAO_Rio_Innovation_Week.pptx — institutional deck for the
"Deep Space and the Lunar Economy" panel (2nd Space Industry Workshop Brazil).

Run:  python3 build/generate_deck.py
Output: dist/MoonDAO_Rio_Innovation_Week.pptx
"""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

import imgutil as I
from deckutil import (
    SLIDE_W, SLIDE_H, MARGIN, CONTENT_TOP, CONTENT_W, FONT, FONT_HEAD,
    NAVY_DARK, NAVY, BLUE, ORANGE, ORANGE_LT, RED, WHITE, BG_LIGHT, BG_PANEL,
    LINE_GRAY, TEXT_DARK, TEXT_GRAY, TEXT_MUTE,
    set_bg, add_rect, add_line, add_text, add_rich, add_bullets, add_picture,
    header, footer, stat_card, section_number_badge, add_donut_chart, logo_lockup_white,
    add_qr_frame, card, accent_bar, accent_rail, evenly_spaced,
)

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.normpath(os.path.join(HERE, '..', 'assets'))
DIST = os.path.normpath(os.path.join(HERE, '..', 'dist'))
os.makedirs(DIST, exist_ok=True)

LOGO = os.path.join(ASSETS, 'MoonDAO_icon.png')
QR = os.path.join(ASSETS, 'qr')
LB = os.path.join(ASSETS, 'leaderboard')
TOTAL_SLIDES = 14

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(BLANK)


# --------------------------------------------------------------------- 01 --
# Title slide
def slide_01():
    s = new_slide()
    set_bg(s, WHITE)

    # Right visual panel — moon + Brazil flag, with a clean orange rail.
    panel_x = Inches(8.5)
    panel_w = Emu(SLIDE_W - panel_x)
    moon_img = I.darken('moon-full.jpg', factor=0.5, ratio=(panel_w.inches, SLIDE_H.inches))
    add_picture(s, moon_img, panel_x, 0, panel_w, SLIDE_H)
    add_rect(s, panel_x, 0, Inches(0.055), SLIDE_H, fill=ORANGE)

    flag = I.rounded_rect_mask('reference/brazil_flag.png', (10, 7), radius_frac=0.05)
    lw_img = panel_w - Inches(1.2)
    lh_img = Emu(int(lw_img * 7 / 10))
    lx_img = panel_x + (panel_w - lw_img) / 2
    ly_img = Inches(4.35)
    add_picture(s, flag, lx_img, ly_img, lw_img, lh_img)
    add_text(s, panel_x + Inches(0.15), ly_img + lh_img + Inches(0.16), panel_w - Inches(0.3), Inches(0.28),
              "Rio de Janeiro, Brazil", size=Pt(12), color=WHITE, bold=True,
              align=PP_ALIGN.CENTER, font=FONT_HEAD)
    add_text(s, panel_x + Inches(0.15), ly_img + lh_img + Inches(0.44), panel_w - Inches(0.3), Inches(0.26),
              "2nd Space Industry Workshop", size=Pt(10), color=RGBColor(0xC7, 0xCC, 0xE4),
              align=PP_ALIGN.CENTER, font=FONT)

    # Left composition — brand first, then title, then speaker.
    lx = Inches(0.85)
    lw = Inches(7.1)
    logo_w = Inches(2.85)
    logo_h = logo_w * (345 / 1233)
    add_picture(s, LOGO, lx, Inches(0.7), logo_w, logo_h)

    add_text(s, lx, Inches(2.05), lw, Inches(0.35),
              "THE INTERNET'S SPACE PROGRAM", size=Pt(13), color=ORANGE, bold=True, font=FONT_HEAD)
    add_text(s, lx - Inches(0.04), Inches(2.4), lw, Inches(1.0),
              "MoonDAO", size=Pt(56), color=NAVY_DARK, bold=True, font=FONT_HEAD)
    add_text(s, lx, Inches(3.45), lw, Inches(0.45),
              "Deep Space and the Lunar Economy", size=Pt(22), color=NAVY, font=FONT_HEAD)
    add_text(s, lx, Inches(3.95), lw, Inches(0.4),
              "2nd Space Industry Workshop Brazil  ·  Brazilian Space Agency (AEB)",
              size=Pt(12.5), color=TEXT_GRAY, font=FONT)

    add_line(s, lx, Inches(4.6), Inches(6.4), 0, color=LINE_GRAY, weight=Pt(0.75))

    circ = I.circle_badge('pablo_headshot.png', size=400)
    cd = Inches(0.9)
    add_picture(s, circ, lx, Inches(4.95), cd, cd)
    add_text(s, lx + cd + Inches(0.22), Inches(5.0), Inches(5.5), Inches(0.35),
              "Pablo Moncada-Larrotiz", size=Pt(17), color=NAVY_DARK, bold=True, font=FONT_HEAD)
    add_text(s, lx + cd + Inches(0.22), Inches(5.38), Inches(5.5), Inches(0.3),
              "Founder & Executive Director, MoonDAO", size=Pt(12.5), color=TEXT_GRAY, font=FONT)

    add_text(s, lx, Inches(6.75), lw, Inches(0.3),
              "pablo@moondao.com   ·   @pablo_moncada_   ·   moondao.com",
              size=Pt(11), color=TEXT_MUTE, font=FONT)
    return s


# --------------------------------------------------------------------- 02 --
def slide_02():
    s = new_slide()
    set_bg(s, WHITE)
    header(s, "About the Speaker", "Pablo Moncada-Larrotiz", 2, LOGO)

    # Left portrait column — centered as a unit.
    left_w = Inches(3.7)
    lx = MARGIN
    badge = I.circle_badge('pablo_headshot.png', size=700)
    bd = Inches(2.4)
    add_picture(s, badge, lx + (left_w - bd) / 2, Inches(1.75), bd, bd)
    add_text(s, lx, Inches(4.35), left_w, Inches(0.35),
              "Pablo Moncada-Larrotiz", size=Pt(16), color=NAVY_DARK, bold=True,
              align=PP_ALIGN.CENTER, font=FONT_HEAD)
    add_text(s, lx, Inches(4.7), left_w, Inches(0.55),
              "Founder & Executive Director\nMoonDAO", size=Pt(12), color=TEXT_GRAY,
              align=PP_ALIGN.CENTER, font=FONT)
    add_text(s, lx, Inches(5.45), left_w, Inches(0.28),
              "San Francisco, CA", size=Pt(11), color=TEXT_MUTE, align=PP_ALIGN.CENTER, font=FONT)
    add_text(s, lx, Inches(5.75), left_w, Inches(0.28),
              "pablo@moondao.com  ·  @pablo_moncada_", size=Pt(10.5), color=TEXT_MUTE,
              align=PP_ALIGN.CENTER, font=FONT)

    # Right content column — bullets + quote, aligned to a shared left edge.
    rx = Inches(5.0)
    rw = SLIDE_W - MARGIN - rx
    bullets = [
        ("Founder of MoonDAO. ", "First elected Executive Lead, since 2021."),
        ("Career: Waymo, YouTube, STEL. ", "Software engineer at Google, then biotech."),
        ("Core contributor, ConstitutionDAO. ", "Helped raise $47M in days for the U.S. Constitution."),
        ("University of Michigan. ", "CS, Mechanical Engineering & Business."),
    ]
    add_bullets(s, rx, Inches(1.85), rw, Inches(3.2), bullets, size=Pt(15), gap=Pt(16))

    card(s, rx, Inches(5.55), rw, Inches(1.05), fill=BG_PANEL, shadow=False, radius=0.1)
    add_text(s, rx + Inches(0.32), Inches(5.7), rw - Inches(0.64), Inches(0.8),
              "\u201cI see decentralization as a way to fix the risks of centralized "
              "control over billions of people's lives.\u201d",
              size=Pt(13), color=NAVY_DARK, italic=True, font=FONT)

    footer(s, 2, TOTAL_SLIDES)
    return s


# --------------------------------------------------------------------- 03 --
def slide_03():
    s = new_slide()
    set_bg(s, WHITE)
    header(s, "Who We Are", "What Is MoonDAO", 3, LOGO)

    add_text(s, MARGIN, Inches(1.55), CONTENT_W, Inches(0.7),
              "An onchain community — anyone, anywhere — pooling resources to fund and "
              "govern humanity's next steps into space.",
              size=Pt(18), color=NAVY_DARK, bold=True, font=FONT_HEAD)

    card(s, MARGIN, Inches(2.45), CONTENT_W, Inches(0.9), fill=BG_PANEL, shadow=False, radius=0.12)
    add_text(s, MARGIN + Inches(0.35), Inches(2.45), Inches(2.4), Inches(0.9),
              "OUR MISSION", size=Pt(11.5), color=ORANGE, bold=True, font=FONT_HEAD, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, MARGIN + Inches(2.8), Inches(2.45), CONTENT_W - Inches(3.15), Inches(0.9),
              "To accelerate the development of a self-sustaining, self-governing settlement on the Moon.",
              size=Pt(15), color=NAVY_DARK, italic=True, font=FONT, anchor=MSO_ANCHOR.MIDDLE)

    bullets = [
        ("Decentralized. ", "Governed onchain via the $MOONEY token."),
        ("Global & permissionless. ", "Anyone, anywhere can join, fund, or build."),
        ("Radically transparent. ", "Every vote and dollar is public onchain."),
    ]
    add_bullets(s, MARGIN, Inches(3.6), CONTENT_W, Inches(1.5), bullets, size=Pt(15), gap=Pt(12))

    stats = [
        ("$8M+", "Raised onchain\nto fund space access"),
        ("12,000+", "$MOONEY\ntoken holders"),
        ("80+", "Space R&D projects\nfunded via governance"),
        ("2", "Crowdfunded citizen-\nastronauts sent to space"),
    ]
    gap = Inches(0.22)
    card_w = evenly_spaced(4, CONTENT_W, gap)
    y0 = Inches(5.4)
    ch = Inches(1.25)
    for i, (num, label) in enumerate(stats):
        cx = MARGIN + i * (card_w + gap)
        stat_card(s, cx, y0, card_w, ch, num, label, accent=[NAVY, ORANGE, RED, BLUE][i % 4])

    footer(s, 3, TOTAL_SLIDES)
    return s


# --------------------------------------------------------------------- 04 --
def slide_04():
    s = new_slide()
    set_bg(s, WHITE)
    header(s, "Our Story", "Five Years of Milestones", 4, LOGO)

    events = [
        ("2021", "MoonDAO Launches", "First DAO focused on space\nexploration & lunar settlement"),
        ("2022", "$8M Raised + First Astronaut", "2,600 ETH raised; Coby Cotton\nflies on Blue Origin's NS-22"),
        ("2023", "80+ Projects Funded", "Community governance funds a\nthriving space R&D ecosystem"),
        ("2024", "Network + 2nd Astronaut", "Space Acceleration Network launches;\nDr. Eiman Jahangir flies to space"),
        ("2025", "Constitution Reaches the Moon", "MoonDAO Launchpad ships; our\nconstitution lands on the lunar surface"),
    ]
    n = len(events)
    card_w = Inches(2.15)
    x0 = MARGIN + card_w / 2
    x1 = SLIDE_W - MARGIN - card_w / 2
    line_y = Inches(2.85)
    add_line(s, x0, line_y, x1 - x0, 0, color=LINE_GRAY, weight=Pt(2.5))
    step = (x1 - x0) / (n - 1)
    colors = [NAVY, BLUE, ORANGE, RED, NAVY]
    for i, (yr, title, desc) in enumerate(events):
        cx = x0 + step * i
        dot = Inches(0.22)
        shp = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - dot / 2, line_y - dot / 2, dot, dot)
        shp.fill.solid()
        shp.fill.fore_color.rgb = colors[i]
        shp.line.color.rgb = WHITE
        shp.line.width = Pt(2)
        shp.shadow.inherit = False

        add_text(s, cx - card_w / 2, line_y - Inches(1.35), card_w, Inches(0.4), yr,
                  size=Pt(20), color=colors[i], bold=True, align=PP_ALIGN.CENTER, font=FONT_HEAD)
        add_text(s, cx - card_w / 2, line_y - Inches(0.98), card_w, Inches(0.7), title,
                  size=Pt(12), color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER, font=FONT_HEAD)

        add_text(s, cx - card_w / 2, line_y + Inches(0.3), card_w, Inches(0.95), desc,
                  size=Pt(10), color=TEXT_GRAY, align=PP_ALIGN.CENTER, font=FONT)

    # Each photo keeps its own native aspect ratio at a shared height, rather
    # than being forced into identical boxes (which stretched the square
    # astronaut portraits horizontally).
    photos = [
        ('MoonDAO-New-Shepard.jpg', 'Team at Blue Origin\u2019s launch site', (16, 10)),
        ('astronaut-coby.png', 'Coby Cotton \u2014 1st MoonDAO astronaut', (1, 1)),
        ('zero-g-image.jpg', 'Zero-gravity training flight', (16, 10)),
        ('astronaut-eiman.png', 'Dr. Eiman Jahangir \u2014 2nd astronaut', (1, 1)),
    ]
    ph = Inches(1.6)
    pgap = Inches(0.28)
    widths = [Emu(int(ph * rw_ / rh_)) for (_, _, (rw_, rh_)) in photos]
    total_pw = sum(widths, Emu(0)) + pgap * (len(photos) - 1)
    px0 = Emu(int((SLIDE_W - total_pw) / 2))
    py0 = Inches(5.0)
    px = px0
    for i, (fname, cap, ratio) in enumerate(photos):
        path = I.rounded_rect_mask(fname, ratio, radius_frac=0.06)
        pw_i = widths[i]
        add_picture(s, path, px, py0, pw_i, ph)
        add_text(s, px - Inches(0.2), py0 + ph + Inches(0.05), pw_i + Inches(0.4), Inches(0.32), cap,
                  size=Pt(9), color=TEXT_GRAY, align=PP_ALIGN.CENTER, font=FONT)
        px = Emu(int(px + pw_i + pgap))

    footer(s, 4, TOTAL_SLIDES)
    return s


# --------------------------------------------------------------------- 05 --
def slide_05():
    s = new_slide()
    set_bg(s, WHITE)
    header(s, "What We Do", "Core Initiatives", 5, LOGO)

    cards = [
        ("Human Spaceflight", "Two crowdfunded citizens sent to space by onchain vote."),
        ("Fund Space R&D", "$600K+ allocated to 80+ projects via open governance."),
        ("Space Training", "Zero-gravity flights and astronaut-prep experiences."),
        ("Space Acceleration Network", "An onchain \u201cstartup society\u201d for the space industry."),
        ("Transparent Governance", "Every vote and treasury move, public onchain."),
        ("Lunar Settlement Roadmap", "Constitution, Launchpad, and missions \u2014 building toward the Moon."),
    ]
    cols = 3
    gap = Inches(0.28)
    cw = evenly_spaced(cols, CONTENT_W, gap)
    ch = Inches(2.05)
    y0 = Inches(1.7)
    accents = [NAVY, ORANGE, BLUE, RED, NAVY, ORANGE]
    for i, (title, desc) in enumerate(cards):
        r, c = divmod(i, cols)
        cx = MARGIN + c * (cw + gap)
        cy = y0 + r * (ch + gap)
        card(s, cx, cy, cw, ch)
        accent_bar(s, cx, cy, cw, accent=accents[i])
        add_text(s, cx + Inches(0.3), cy + Inches(0.32), cw - Inches(0.6), Inches(0.5),
                  title, size=Pt(15), color=NAVY_DARK, bold=True, font=FONT_HEAD)
        add_text(s, cx + Inches(0.3), cy + Inches(0.9), cw - Inches(0.6), Inches(0.95),
                  desc, size=Pt(12), color=TEXT_GRAY, font=FONT)

    footer(s, 5, TOTAL_SLIDES)
    return s


# --------------------------------------------------------------------- 06 --
def slide_06():
    s = new_slide()
    set_bg(s, NAVY_DARK)
    bg = I.darken('reference/earth_sunrise.jpeg', factor=0.42, ratio=(SLIDE_W.inches, SLIDE_H.inches))
    add_picture(s, bg, 0, 0, SLIDE_W, SLIDE_H)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.09), fill=ORANGE)

    add_text(s, Inches(1.1), Inches(2.55), Inches(11.1), Inches(0.45),
              "THE NEXT OPPORTUNITY", size=Pt(15), color=ORANGE, bold=True,
              align=PP_ALIGN.CENTER, font=FONT_HEAD)
    add_text(s, Inches(1.1), Inches(3.05), Inches(11.1), Inches(1.5),
              "Go to Space with Frank White", size=Pt(46), color=WHITE, bold=True,
              align=PP_ALIGN.CENTER, font=FONT_HEAD)
    add_text(s, Inches(1.8), Inches(4.15), Inches(9.7), Inches(0.7),
              "Funding a flight for the author of \u201cThe Overview Effect\u201d \u2014 and one "
              "MoonDAO community member \u2014 to finally see Earth from space.",
              size=Pt(15.5), color=RGBColor(0xE6, 0xE9, 0xF5), align=PP_ALIGN.CENTER, font=FONT)

    footer_light(s, 6)
    return s


def footer_light(s, page_no):
    add_text(s, MARGIN, SLIDE_H - Inches(0.42), Inches(9.5), Inches(0.28),
              "MoonDAO  ·  Deep Space and the Lunar Economy  ·  2nd Space Industry Workshop Brazil",
              size=Pt(9), color=RGBColor(0xAF, 0xB6, 0xCC), font=FONT)
    add_text(s, SLIDE_W - MARGIN - Inches(1.3), SLIDE_H - Inches(0.42), Inches(1.3), Inches(0.28),
              f"{page_no:02d}  /  {TOTAL_SLIDES}", size=Pt(9), color=RGBColor(0xAF, 0xB6, 0xCC),
              font=FONT_HEAD, align=PP_ALIGN.RIGHT)


# --------------------------------------------------------------------- 07 --
def slide_07():
    s = new_slide()
    set_bg(s, WHITE)
    header(s, "The Next Opportunity", "Who Is Frank White", 7, LOGO)

    lx, lw = MARGIN, Inches(3.9)
    badge = I.circle_badge('reference/frank_white_real_photo.png', size=700)
    bd = Inches(2.3)
    add_picture(s, badge, lx + (lw - bd) / 2, Inches(1.7), bd, bd)
    add_text(s, lx, Inches(4.15), lw, Inches(0.4), "Frank White", size=Pt(16),
              color=NAVY_DARK, bold=True, align=PP_ALIGN.CENTER, font=FONT_HEAD)
    add_text(s, lx, Inches(4.5), lw, Inches(0.4), "Author, \u201cThe Overview Effect\u201d",
              size=Pt(12), color=TEXT_GRAY, align=PP_ALIGN.CENTER, font=FONT)

    book = I.rounded_rect_mask('reference/overview_effect_book.jpg', (2, 3), radius_frac=0.03)
    book_h = Inches(1.95)
    book_w = Emu(int(book_h * 2 / 3))
    add_picture(s, book, lx + (lw - book_w) / 2, Inches(4.95), book_w, book_h)
    add_text(s, lx, Inches(4.95) + book_h + Inches(0.06), lw, Inches(0.3),
              "\u201cThe Overview Effect,\u201d 4th ed.", size=Pt(9), color=TEXT_MUTE,
              align=PP_ALIGN.CENTER, font=FONT)

    rx = Inches(5.15)
    rw = SLIDE_W - MARGIN - rx
    bullets = [
        ("Coined \u201cthe Overview Effect.\u201d ", "The name for the shift astronauts feel seeing Earth from space."),
        ("40 years of research. ", "Interviewed 50+ astronauts and cosmonauts."),
        ("Shaped the industry. ", "Co-founded the Overview Institute; inspired a documentary."),
        ("Has never been to space.", ""),
    ]
    add_bullets(s, rx, Inches(1.85), rw, Inches(3.0), bullets, size=Pt(16), gap=Pt(18))

    card(s, rx, Inches(5.2), rw, Inches(1.25), fill=BG_PANEL, shadow=False, radius=0.1)
    add_text(s, rx + Inches(0.3), Inches(5.35), rw - Inches(0.6), Inches(1.0),
              "\u201cAstronauts describe seeing Earth from space as an instant, visceral realization "
              "that we share one planet, one atmosphere, one destiny \u2014 the Overview Effect.\u201d",
              size=Pt(12.5), color=NAVY_DARK, italic=True, font=FONT)

    footer(s, 7, TOTAL_SLIDES)
    return s


# --------------------------------------------------------------------- 08 --
def slide_08():
    s = new_slide()
    set_bg(s, WHITE)
    header(s, "The Next Opportunity", "Fly to Space with Frank White", 8, LOGO)

    # Three equal columns for a more symmetrical composition.
    gap = Inches(0.35)
    col_w = evenly_spaced(3, CONTENT_W, gap)
    y0 = Inches(1.6)

    # Left — funding progress card.
    lx = MARGIN
    card(s, lx, y0, col_w, Inches(4.35), fill=BG_LIGHT, shadow=False)
    raised, goal = 172, 250
    pct = round(raised / goal * 100)
    add_donut_chart(s, lx + Inches(0.45), y0 + Inches(0.35), col_w - Inches(0.9), Inches(2.35),
                      ["Raised", "Remaining"], [raised, goal - raised],
                      [ORANGE, LINE_GRAY], hole_size=68)
    add_text(s, lx + Inches(0.45), y0 + Inches(1.15), col_w - Inches(0.9), Inches(0.7),
              f"{pct}%\nfunded", size=Pt(16), color=NAVY_DARK, bold=True,
              align=PP_ALIGN.CENTER, font=FONT_HEAD)
    add_text(s, lx + Inches(0.2), y0 + Inches(2.9), col_w - Inches(0.4), Inches(0.4),
              f"${raised}K of ${goal}K", size=Pt(18), color=NAVY_DARK, bold=True,
              align=PP_ALIGN.CENTER, font=FONT_HEAD)
    add_text(s, lx + Inches(0.2), y0 + Inches(3.4), col_w - Inches(0.4), Inches(0.6),
              "157 contributors\ngoal secures a seat", size=Pt(11.5), color=TEXT_GRAY,
              align=PP_ALIGN.CENTER, font=FONT)

    # Middle — top 3 candidates.
    mx = lx + col_w + gap
    card(s, mx, y0, col_w, Inches(4.35), fill=BG_LIGHT, shadow=False)
    add_text(s, mx + Inches(0.25), y0 + Inches(0.22), col_w - Inches(0.5), Inches(0.35),
              "Top candidates", size=Pt(13), color=NAVY_DARK, bold=True, font=FONT_HEAD)
    add_text(s, mx + Inches(0.25), y0 + Inches(0.52), col_w - Inches(0.5), Inches(0.28),
              "25 advance to Round 2", size=Pt(11), color=TEXT_MUTE, font=FONT)
    leaderboard = [
        (1, 'Andrew "Titan" Parris', 2083, os.path.join(LB, 'citizen_7_parris.png'), ORANGE),
        (2, 'EngiBob', 680, os.path.join(LB, 'citizen_engibob.png'), NAVY),
        (3, 'Jas', 652, os.path.join(LB, 'citizen_184_jas.png'), BLUE),
    ]
    ry = y0 + Inches(1.0)
    row_h = Inches(0.95)
    pd = Inches(0.55)
    for rank, name, amt, photo, color in leaderboard:
        add_text(s, mx + Inches(0.22), ry + Inches(0.12), Inches(0.3), Inches(0.3), f"{rank}",
                  size=Pt(14), color=TEXT_MUTE, bold=True, font=FONT_HEAD)
        badge = I.circle_badge(os.path.relpath(photo, ASSETS), size=300)
        add_picture(s, badge, mx + Inches(0.55), ry + Inches(0.08), pd, pd)
        add_text(s, mx + Inches(1.25), ry + Inches(0.05), col_w - Inches(1.55), Inches(0.28),
                  name, size=Pt(12), color=NAVY_DARK, bold=True, font=FONT_HEAD)
        bar_max = col_w - Inches(1.55)
        add_rect(s, mx + Inches(1.25), ry + Inches(0.38), bar_max, Inches(0.12), fill=WHITE,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        add_rect(s, mx + Inches(1.25), ry + Inches(0.38), Emu(int(bar_max * (amt / 2083))), Inches(0.12),
                  fill=color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        add_text(s, mx + Inches(1.25), ry + Inches(0.55), col_w - Inches(1.55), Inches(0.24),
                  f"{amt:,} $OVERVIEW", size=Pt(10), color=TEXT_GRAY, font=FONT)
        ry += row_h
    add_text(s, mx + Inches(0.25), y0 + Inches(3.95), col_w - Inches(0.5), Inches(0.28),
              "Any Citizen can enter and climb.", size=Pt(10), color=TEXT_MUTE,
              italic=True, font=FONT)

    # Right — QR code column.
    rx = mx + col_w + gap
    card(s, rx, y0, col_w, Inches(4.35), fill=BG_LIGHT, shadow=False)
    qr_size = Inches(2.0)
    qr_path = os.path.join(QR, 'qr_leaderboard.png')
    qlx = rx + (col_w - qr_size - Inches(0.36)) / 2
    fw, fh = add_qr_frame(s, qr_path, qlx, y0 + Inches(0.55), qr_size)
    add_text(s, rx + Inches(0.2), y0 + Inches(0.55) + fh + Inches(0.25), col_w - Inches(0.4), Inches(0.55),
              "Support Frank &\nBack a Candidate", size=Pt(13), color=NAVY_DARK, bold=True,
              align=PP_ALIGN.CENTER, font=FONT_HEAD)
    add_text(s, rx + Inches(0.2), y0 + Inches(0.55) + fh + Inches(0.85), col_w - Inches(0.4), Inches(0.28),
              "moondao.com/overview-vote", size=Pt(10), color=TEXT_MUTE,
              align=PP_ALIGN.CENTER, font=FONT)

    card(s, MARGIN, Inches(6.15), CONTENT_W, Inches(0.5), fill=BG_PANEL, shadow=False, radius=0.2)
    add_text(s, MARGIN + Inches(0.35), Inches(6.15), CONTENT_W - Inches(0.7), Inches(0.5),
              "$100+ grants free citizenship \u2014 and a shot at flying alongside Frank.",
              size=Pt(12.5), color=NAVY_DARK, italic=True, font=FONT, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, 8, TOTAL_SLIDES)
    return s


# --------------------------------------------------------------------- 09 --
def slide_09():
    s = new_slide()
    set_bg(s, WHITE)
    header(s, "Funding Innovation", "DePrize: A Prediction Market for Space Delivery", 9, LOGO,
            title_size=Pt(24))

    # Mechanism — three equal steps.
    steps = [
        ("1", "Back a Team", "Wager ETH on who delivers first", NAVY),
        ("2", "Live Odds", "Every bet grows the shared prize", BLUE),
        ("3", "Winner Delivers", "Community declares winner; backers get paid", ORANGE),
    ]
    gap = Inches(0.28)
    cw = evenly_spaced(3, CONTENT_W, gap)
    y0 = Inches(1.55)
    ch = Inches(1.25)
    for i, (num, title, desc, accent) in enumerate(steps):
        cx = MARGIN + i * (cw + gap)
        card(s, cx, y0, cw, ch, fill=BG_LIGHT, shadow=False)
        section_number_badge(s, cx + Inches(0.22), y0 + Inches(0.22), num, accent=accent, d=Inches(0.42))
        add_text(s, cx + Inches(0.78), y0 + Inches(0.22), cw - Inches(1.0), Inches(0.42),
                  title, size=Pt(14.5), color=NAVY_DARK, bold=True, font=FONT_HEAD, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cx + Inches(0.22), y0 + Inches(0.78), cw - Inches(0.44), Inches(0.35),
                  desc, size=Pt(11), color=TEXT_GRAY, font=FONT)

    # Tech tree verticals — 2×4 chip grid.
    vy0 = y0 + ch + Inches(0.28)
    add_text(s, MARGIN, vy0, CONTENT_W, Inches(0.28),
              "8 CAPABILITY VERTICALS ON THE TECH TREE", size=Pt(11), color=ORANGE, bold=True, font=FONT_HEAD)
    verticals = [
        ("Comms / PNT", NAVY), ("Surface Construction", ORANGE), ("Habitat", BLUE), ("ISRU Plant", RED),
        ("Power", NAVY), ("Rover", ORANGE), ("Crewed Base", BLUE), ("Lander", RED),
    ]
    vcols = 4
    vgap = Inches(0.16)
    vcw = evenly_spaced(vcols, CONTENT_W, vgap)
    vch = Inches(0.42)
    vy = vy0 + Inches(0.35)
    for i, (name, accent) in enumerate(verticals):
        r, c = divmod(i, vcols)
        vx = MARGIN + c * (vcw + vgap)
        cy = vy + r * (vch + Inches(0.12))
        card(s, vx, cy, vcw, vch, fill=BG_LIGHT, shadow=False, radius=0.4)
        dot = Inches(0.1)
        dot_shp = s.shapes.add_shape(MSO_SHAPE.OVAL, vx + Inches(0.18), cy + (vch - dot) / 2, dot, dot)
        dot_shp.fill.solid()
        dot_shp.fill.fore_color.rgb = accent
        dot_shp.line.fill.background()
        dot_shp.shadow.inherit = False
        add_text(s, vx + Inches(0.4), cy, vcw - Inches(0.55), vch, name,
                  size=Pt(11), color=NAVY_DARK, bold=True, font=FONT_HEAD, anchor=MSO_ANCHOR.MIDDLE)

    # Worked math — bottom band.
    my0 = vy + 2 * (vch + Inches(0.12)) + Inches(0.22)
    mh = Inches(1.85)
    card(s, MARGIN, my0, CONTENT_W, mh, fill=BG_PANEL, shadow=False)
    add_text(s, MARGIN + Inches(0.35), my0 + Inches(0.18), Inches(4), Inches(0.28),
              "THE MATH, WORKED", size=Pt(11), color=ORANGE, bold=True, font=FONT_HEAD)

    steps_math = [
        "Bet 1 ETH on a provider priced at 40% odds",
        "5% \u2192 prize pool  ·  95% \u2192 market (2.375 shares at $0.40)",
        "Less the 1% LMSR fee \u2192 \u2248 2.35 shares held",
        "If they win: shares redeem 1:1 \u2192 \u2248 2.32 ETH back",
    ]
    ty = my0 + Inches(0.5)
    for line in steps_math:
        add_text(s, MARGIN + Inches(0.35), ty, Inches(8.0), Inches(0.28), f"\u2192  {line}",
                  size=Pt(11.5), color=TEXT_DARK, font=FONT)
        ty += Inches(0.28)

    divider_x = MARGIN + Inches(8.7)
    add_line(s, divider_x, my0 + Inches(0.25), 0, mh - Inches(0.5), color=LINE_GRAY, weight=Pt(0.75))
    rmx = divider_x + Inches(0.25)
    rmw = SLIDE_W - MARGIN - rmx - Inches(0.15)
    add_text(s, rmx, my0 + Inches(0.4), rmw, Inches(0.55), "+1.32 ETH", size=Pt(24),
              color=ORANGE, bold=True, align=PP_ALIGN.CENTER, font=FONT_HEAD)
    add_text(s, rmx, my0 + Inches(1.0), rmw, Inches(0.6),
              "net gain + 50 $OVERVIEW,\nfunded by other bettors",
              size=Pt(10.5), color=TEXT_GRAY, align=PP_ALIGN.CENTER, font=FONT)

    footer(s, 9, TOTAL_SLIDES)
    return s


# --------------------------------------------------------------------- 10 --
def slide_10():
    s = new_slide()
    set_bg(s, WHITE)
    header(s, "Funding Innovation", "Surviving the 354-Hour Lunar Night", 10, LOGO,
            title_size=Pt(26))

    cols = [
        ("THE PHYSICS", "354 hrs", "of total darkness, every cycle", NAVY, [
            "The Moon is tidally locked: one rotation = one 29.5-day orbit",
            "\u2192 sunrise to sunrise (\u201clunar day\u201d) = 29.5 Earth days",
            "\u2192 continuous night \u2248 14.75 days = 354 hours",
            "No atmosphere \u2192 no convection \u2014 heat only radiates away",
        ]),
        ("THE EXTREMES", "257\u00b0C", "swing between sun and shadow", RED, [
            "South Pole sunlit ground: up to 54\u00b0C (130\u00b0F)",
            "Nearby permanent shadow: as low as \u2212203\u00b0C (\u2212334\u00b0F)",
            "Colder than anywhere ever recorded on Earth",
            "Regolith insulates the subsurface, but can't power a base",
        ]),
        ("THE MATH", "120 tons", "of batteries \u2014 for one night", ORANGE, [
            "40 kWe load \u00d7 354 hrs \u2248 14,160 kWh needed",
            "At Tesla Powerwall density (13.5 kWh / 114 kg):",
            "\u2192 \u2248 1,050 units \u2248 120 metric tons of batteries",
            "NASA's whole heavy-lander budget: 15 tons, total",
        ]),
    ]
    gap = Inches(0.28)
    cw = evenly_spaced(3, CONTENT_W, gap)
    y0 = Inches(1.55)
    ch = Inches(4.05)
    for i, (label, stat, sub, accent, lines) in enumerate(cols):
        cx = MARGIN + i * (cw + gap)
        card(s, cx, y0, cw, ch, fill=BG_LIGHT, shadow=False)
        accent_bar(s, cx, y0, cw, accent=accent)
        add_text(s, cx + Inches(0.26), y0 + Inches(0.22), cw - Inches(0.52), Inches(0.28),
                  label, size=Pt(10.5), color=accent, bold=True, font=FONT_HEAD)
        add_text(s, cx + Inches(0.26), y0 + Inches(0.52), cw - Inches(0.52), Inches(0.55),
                  stat, size=Pt(26), color=NAVY_DARK, bold=True, font=FONT_HEAD)
        add_text(s, cx + Inches(0.26), y0 + Inches(1.1), cw - Inches(0.52), Inches(0.4),
                  sub, size=Pt(10.5), color=TEXT_GRAY, italic=True, font=FONT)
        add_line(s, cx + Inches(0.26), y0 + Inches(1.6), cw - Inches(0.52), 0, color=LINE_GRAY, weight=Pt(0.75))
        ly = y0 + Inches(1.8)
        for line in lines:
            add_text(s, cx + Inches(0.26), ly, cw - Inches(0.52), Inches(0.55), f"\u2022  {line}",
                      size=Pt(10), color=TEXT_DARK, font=FONT)
            ly += Inches(0.55)

    card(s, MARGIN, Inches(5.8), CONTENT_W, Inches(0.62), fill=NAVY_DARK, shadow=False, radius=0.12)
    add_text(s, MARGIN + Inches(0.35), Inches(5.8), CONTENT_W - Inches(0.7), Inches(0.62),
              "THE ANSWER \u2014 a 40 kWe fission reactor, shielding, and radiators fit the same 15-ton "
              "lander and run for a 10-year design life. NASA calls this civil-space shortfall #1.",
              size=Pt(11.5), color=WHITE, italic=True, font=FONT, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, MARGIN, Inches(6.55), CONTENT_W, Inches(0.28),
              "Sources: NASA \u2014 The Harsh Environment of the Lunar South Pole; NASA Fission Surface "
              "Power Program; NASA HLS Lunar Thermal Analysis Guidebook.",
              size=Pt(8), color=TEXT_MUTE, italic=True, font=FONT)

    footer(s, 10, TOTAL_SLIDES)
    return s


# --------------------------------------------------------------------- 11 --
def slide_11():
    s = new_slide()
    set_bg(s, NAVY_DARK)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.09), fill=ORANGE)

    add_text(s, MARGIN, Inches(0.32), Inches(10), Inches(0.3),
              "FUNDING INNOVATION", size=Pt(12.5), color=ORANGE, bold=True, font=FONT_HEAD)
    add_text(s, MARGIN, Inches(0.6), Inches(9), Inches(0.5),
              "Moonbase Zero", size=Pt(26), color=WHITE, bold=True, font=FONT_HEAD)
    WHITE_MARK = os.path.join(ASSETS, 'moondao_mark_white.png')
    mark_h = Inches(0.55)
    mark_w = mark_h * (620 / 650)
    add_picture(s, WHITE_MARK, SLIDE_W - MARGIN - mark_w, Inches(0.4), mark_w, mark_h)

    # The real, live Moonbase Zero screenshot — captured directly from a running
    # instance of ui/pages/moonbase/index.tsx (see README for how this was done).
    RATIO = (1300, 946)
    shot = I.rounded_rect_mask('moonbase_real/moonbase_main.png', RATIO, radius_frac=0.025)
    sh = Inches(4.55)
    sw = Emu(int(sh * RATIO[0] / RATIO[1]))
    sx = MARGIN
    sy = Inches(1.3)
    add_picture(s, shot, sx, sy, sw, sh)

    add_rect(s, sx, sy + sh + Inches(0.14), Inches(1.95), Inches(0.38),
              fill=ORANGE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    add_text(s, sx, sy + sh + Inches(0.14), Inches(1.95), Inches(0.38),
              "IN DEVELOPMENT", size=Pt(10), color=WHITE, bold=True, align=PP_ALIGN.CENTER,
              font=FONT_HEAD, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, sx + Inches(2.15), sy + sh + Inches(0.14), Inches(3), Inches(0.38),
              "moondao.com/moonbase", size=Pt(11), color=RGBColor(0xC7, 0xCC, 0xE4), font=FONT,
              anchor=MSO_ANCHOR.MIDDLE)

    # Secondary real screenshot — a capability-race detail panel, opened live.
    dw = SLIDE_W - MARGIN - (sx + sw + Inches(0.35))
    dh = Emu(int(dw * RATIO[1] / RATIO[0]))
    dx = sx + sw + Inches(0.35)
    dy = sy
    detail = I.rounded_rect_mask('moonbase_real/moonbase_lander_detail.png', RATIO, radius_frac=0.025)
    add_picture(s, detail, dx, dy, dw, dh)
    add_text(s, dx, dy + dh + Inches(0.12), dw, Inches(0.3),
              "Click a race to see live competitors, criteria, and sources",
              size=Pt(10.5), color=RGBColor(0xC7, 0xCC, 0xE4), italic=True, font=FONT,
              align=PP_ALIGN.CENTER)

    add_text(s, MARGIN, Inches(6.6), SLIDE_W - 2 * MARGIN, Inches(0.35),
              "A true-to-scale 3D site on the Shackleton connecting ridge \u2014 8 capability races, "
              "24 real competitors, all publicly sourced.",
              size=Pt(10.5), color=RGBColor(0x9A, 0xA1, 0xBF), italic=True, font=FONT,
              align=PP_ALIGN.CENTER)

    footer_light(s, 11)
    return s


# --------------------------------------------------------------------- 12 --
def slide_12():
    s = new_slide()
    set_bg(s, WHITE)
    header(s, "The Bigger Picture", "A Level Playing Field for the Next Space Age", 12, LOGO,
            title_size=Pt(27))

    lx, lw = MARGIN, Inches(7.1)
    add_text(s, lx, Inches(1.7), lw, Inches(0.85),
              "Deep space used to require a national agency \u2014 or a billionaire.",
              size=Pt(18), color=NAVY_DARK, bold=True, font=FONT_HEAD)

    bullets = [
        ("Permissionless. ", "Any country, company, or individual can fund and compete."),
        ("Already global. ", "25+ countries in the Space Acceleration Network."),
        ("An open invitation. ", "To anyone in this room \u2014 and your own ecosystem back home."),
    ]
    add_bullets(s, lx, Inches(2.7), lw, Inches(2.6), bullets, size=Pt(16.5), gap=Pt(20))

    # The real Space Acceleration Network graphic, shown in full (native ratio, no crop).
    san = I.rounded_rect_mask('reference/space_acceleration_network.png', (1600, 902), radius_frac=0.04)
    rw_img = SLIDE_W - MARGIN - (lx + lw + Inches(0.4))
    rh_img = Emu(int(rw_img * 902 / 1600))
    rx_img = lx + lw + Inches(0.4)
    ry_img = Inches(2.65)
    add_picture(s, san, rx_img, ry_img, rw_img, rh_img)

    stats = [
        ("25+", "Countries in the\nSpace Acceleration Network"),
        ("12,000+", "$MOONEY\ntoken holders worldwide"),
        ("100%", "Onchain, transparent\nvotes & treasury"),
    ]
    gap = Inches(0.28)
    card_w = evenly_spaced(3, CONTENT_W, gap)
    y0 = Inches(5.7)
    ch = Inches(1.05)
    for i, (num, label) in enumerate(stats):
        cx = MARGIN + i * (card_w + gap)
        stat_card(s, cx, y0, card_w, ch, num, label, accent=[NAVY, ORANGE, RED][i])

    footer(s, 12, TOTAL_SLIDES)
    return s


# --------------------------------------------------------------------- 13 --
def slide_13():
    s = new_slide()
    set_bg(s, NAVY_DARK)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.09), fill=ORANGE)

    add_text(s, MARGIN, Inches(0.32), Inches(10), Inches(0.3),
              "THE BIGGER PICTURE", size=Pt(12.5), color=ORANGE, bold=True, font=FONT_HEAD)
    add_text(s, MARGIN, Inches(0.6), Inches(10), Inches(0.5),
              "Our Global Citizens", size=Pt(26), color=WHITE, bold=True, font=FONT_HEAD)
    WHITE_MARK = os.path.join(ASSETS, 'moondao_mark_white.png')
    mark_h = Inches(0.55)
    mark_w = mark_h * (620 / 650)
    add_picture(s, WHITE_MARK, SLIDE_W - MARGIN - mark_w, Inches(0.4), mark_w, mark_h)

    # The real MoonDAO Network / Citizens page graphic, in full — fit by height so it
    # never overflows the slide, then centered horizontally.
    RATIO2 = (1600, 900)
    net = I.rounded_rect_mask('reference/moondao_citizens_network.png', RATIO2, radius_frac=0.02)
    nh = Inches(5.15)
    nw = Emu(int(nh * RATIO2[0] / RATIO2[1]))
    nx = (SLIDE_W - nw) / 2
    ny = Inches(1.3)
    add_picture(s, net, nx, ny, nw, nh)

    add_text(s, MARGIN, ny + nh + Inches(0.14), SLIDE_W - 2 * MARGIN, Inches(0.35),
              "12,000+ $MOONEY holders and a growing roster of onchain \u201ccitizens\u201d \u2014 engineers, "
              "artists, and space professionals \u2014 building this together.",
              size=Pt(11), color=RGBColor(0x9A, 0xA1, 0xBF), italic=True, font=FONT,
              align=PP_ALIGN.CENTER)

    footer_light(s, 13)
    return s


# --------------------------------------------------------------------- 14 --
def slide_14():
    s = new_slide()
    set_bg(s, NAVY_DARK)
    panel_img = I.darken('earthrise.jpg', factor=0.28, ratio=(SLIDE_W.inches, SLIDE_H.inches))
    add_picture(s, panel_img, 0, 0, SLIDE_W, SLIDE_H)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.09), fill=ORANGE)

    WHITE_MARK = os.path.join(ASSETS, 'moondao_mark_white.png')
    lockup_h = Inches(0.62)
    lockup_w = lockup_h * (620 / 650) + Inches(0.12) + Inches(1.55)
    logo_lockup_white(s, WHITE_MARK, (SLIDE_W - lockup_w) / 2, Inches(0.8), lockup_h)

    add_text(s, Inches(1.2), Inches(1.65), Inches(10.9), Inches(0.7),
              "Let's Build the Lunar Economy Together", size=Pt(30), color=WHITE, bold=True,
              align=PP_ALIGN.CENTER, font=FONT_HEAD)

    # Left: speaker. Right: QR to join.
    badge = I.circle_badge('pablo_headshot.png', size=500, border_color=(255, 255, 255, 255), border_px=10)
    bd = Inches(1.4)
    lcx = Inches(3.7)
    add_picture(s, badge, lcx - bd / 2, Inches(2.75), bd, bd)
    add_text(s, lcx - Inches(2.5), Inches(4.35), Inches(5), Inches(0.4),
              "Pablo Moncada-Larrotiz", size=Pt(15), color=WHITE, bold=True,
              align=PP_ALIGN.CENTER, font=FONT_HEAD)
    add_text(s, lcx - Inches(2.5), Inches(4.72), Inches(5), Inches(0.35),
              "Founder & Executive Director, MoonDAO", size=Pt(11.5),
              color=RGBColor(0xCB, 0xD1, 0xE6), align=PP_ALIGN.CENTER, font=FONT)
    add_text(s, lcx - Inches(2.5), Inches(5.1), Inches(5), Inches(0.35),
              "pablo@moondao.com   \u00b7   @pablo_moncada_", size=Pt(11.5),
              color=RGBColor(0xCB, 0xD1, 0xE6), align=PP_ALIGN.CENTER, font=FONT)

    add_line(s, Inches(6.9), Inches(2.8), 0, Inches(2.85), color=RGBColor(0x3A, 0x44, 0x70), weight=Pt(1))

    qr_size = Inches(1.75)
    qr_path = os.path.join(QR, 'qr_join_moondao.png')
    qcx = Inches(9.6)
    fw, fh = add_qr_frame(s, qr_path, qcx - qr_size / 2 - Inches(0.16), Inches(2.75), qr_size,
                            frame_fill=WHITE, line=WHITE)
    add_text(s, qcx - Inches(1.7), Inches(2.75) + fh + Inches(0.14), Inches(3.4), Inches(0.35),
              "Scan to Join MoonDAO", size=Pt(13.5), color=WHITE, bold=True,
              align=PP_ALIGN.CENTER, font=FONT_HEAD)
    add_text(s, qcx - Inches(1.7), Inches(2.75) + fh + Inches(0.5), Inches(3.4), Inches(0.3),
              "moondao.com/join", size=Pt(10.5), color=RGBColor(0xCB, 0xD1, 0xE6),
              align=PP_ALIGN.CENTER, font=FONT)

    ctas = ["moondao.com", "Discord: moondao.com/discord"]
    add_text(s, Inches(1.2), Inches(6.55), Inches(10.9), Inches(0.4),
              "   \u00b7   ".join(ctas), size=Pt(12.5), color=ORANGE_LT, bold=True,
              align=PP_ALIGN.CENTER, font=FONT_HEAD)

    footer_light(s, 14)
    return s


# -------------------------------------------------------- speaker notes --
# Target: ~8 minutes total (the panel's suggested opening-remarks slot).
NOTES = {
    1: "[~20s] Quick intro — who you are and why MoonDAO is on this panel. Don't linger, the "
       "content does the introducing.",
    2: "[~30s] One or two sentences max — the room can read the bullets. Land on the ConstitutionDAO "
       "line, it's the most relatable proof point for a non-crypto audience.",
    3: "[~40s] Define MoonDAO in one breath before the stats: an onchain community funding and "
       "governing space access. Let the four numbers do the heavy lifting.",
    4: "[~45s] Walk left to right once, quickly. Slow down only on 2022 (first astronaut) and "
       "2025 (constitution on the Moon) — those are the two moments that make people lean in.",
    5: "[~40s] Don't read every card. Group them verbally: 'flight access, funding, training, "
       "network, governance, and the roadmap that ties it together.'",
    6: "[~10s] Let this breathe — pause on the title before moving to Frank White. This is the "
       "pivot from 'here's who we are' to 'here's the live opportunity.'",
    7: "[~45s] Frank is the emotional core of this section — 40 years, 50+ astronauts, coined the "
       "term, and has never flown himself. That irony is the hook.",
    8: "[~40s] Point at the QR code directly — invite the room to scan right now while you talk. "
       "These are real, live candidates on the board; anyone can enter and climb it.",
    9: "[~50s] Walk the 3-step mechanism fast, name the 8 verticals as one group ('comms, "
       "construction, habitat, ISRU, power, rover, crewed base, lander'), then slow down for the "
       "worked math — the +1.32 ETH number is the one that lands. Frame it as a new financing "
       "primitive, not 'crypto betting.'",
    10: "[~55s] This is the technical-depth moment. Walk physics \u2192 extremes \u2192 math in order, "
        "landing on '120 tons of batteries vs. a 15-ton lander.' It shows MoonDAO understands the "
        "engineering, not just the funding mechanics.",
    11: "[~35s] This is the real, live product — say so explicitly. Point out the second screenshot "
        "shows an actual capability race opened live, with real competitors and sources.",
    12: "[~40s] This is the slide that answers the panel's brief directly — make it personal to "
        "whoever's in the room, not just the named institutions.",
    13: "[~15s] Quick, human beat before the close — real people, real names, behind the "
        "onchain numbers. Don't over-explain, let the faces land.",
    14: "[~20s] Point at the QR code — invite people to scan and join MoonDAO before you leave the "
        "stage. Contact info stays up during Q&A.",
}


# ------------------------------------------------------------------ build --
for i, fn in enumerate([slide_01, slide_02, slide_03, slide_04, slide_05, slide_06,
                          slide_07, slide_08, slide_09, slide_10, slide_11, slide_12,
                          slide_13, slide_14], start=1):
    slide = fn()
    note = NOTES.get(i)
    if note:
        slide.notes_slide.notes_text_frame.text = note

out_path = os.path.join(DIST, 'MoonDAO_Rio_Innovation_Week.pptx')
prs.save(out_path)
print('Saved:', out_path)
