"""Shared drawing helpers for the MoonDAO / Rio Innovation Week deck.

Design system: neutral, institutional. White/light-gray canvas, MoonDAO's own
brand navy / orange / red used sparingly as accents, one sans-serif family
throughout, consistent header + footer treatment, generous whitespace.
"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import copy

# ---------------------------------------------------------------- palette --
NAVY_DARK = RGBColor(0x1A, 0x22, 0x4A)
NAVY = RGBColor(0x22, 0x35, 0x7A)
BLUE = RGBColor(0x2B, 0x4C, 0x9B)
ORANGE = RGBColor(0xE8, 0x8B, 0x1F)
ORANGE_LT = RGBColor(0xF6, 0xB9, 0x62)
RED = RGBColor(0xA3, 0x2E, 0x35)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF6, 0xF7, 0xFB)
BG_PANEL = RGBColor(0xEE, 0xF0, 0xF7)
LINE_GRAY = RGBColor(0xDD, 0xE1, 0xEC)
TEXT_DARK = RGBColor(0x20, 0x24, 0x33)
TEXT_GRAY = RGBColor(0x5B, 0x62, 0x78)
TEXT_MUTE = RGBColor(0x8A, 0x90, 0xA3)

FONT = "Calibri"
FONT_HEAD = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN = Inches(0.55)


def set_bg(slide, color=WHITE):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _no_line(shape):
    shape.line.fill.background()


def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0.75),
             shadow=False, shape_type=MSO_SHAPE.RECTANGLE, radius=None):
    shp = slide.shapes.add_shape(shape_type, l, t, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        _no_line(shp)
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w
    shp.shadow.inherit = False
    if shadow:
        el = shp._element.spPr
        effectLst = el.makeelement(qn('a:effectLst'), {})
        outerShdw = el.makeelement(qn('a:outerShdw'), {
            'blurRad': '90000', 'dist': '38000', 'dir': '5400000', 'rotWithShape': '0'
        })
        clr = el.makeelement(qn('a:srgbClr'), {'val': '1A224A'})
        alpha = el.makeelement(qn('a:alpha'), {'val': '18000'})
        clr.append(alpha)
        outerShdw.append(clr)
        effectLst.append(outerShdw)
        el.append(effectLst)
    if radius is not None and shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    return shp


def add_line(slide, l, t, w, h, color=LINE_GRAY, weight=Pt(1)):
    ln = slide.shapes.add_connector(1, l, t, l + w, t + h)
    ln.line.color.rgb = color
    ln.line.width = weight
    return ln


def _set_paragraph(p, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT,
                    font=FONT, spacing=None, space_after=None, space_before=None):
    p.text = text
    p.alignment = align
    if spacing is not None:
        p.line_spacing = spacing
    if space_after is not None:
        p.space_after = space_after
    if space_before is not None:
        p.space_before = space_before
    for run in p.runs:
        run.font.size = size
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
        run.font.color.rgb = color


def add_text(slide, l, t, w, h, text, size=Pt(18), color=TEXT_DARK, bold=False,
             italic=False, align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP,
             spacing=None, wrap=True, autosize=False, shrink=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split('\n')
    p = tf.paragraphs[0]
    _set_paragraph(p, lines[0], size, color, bold, italic, align, font, spacing)
    for extra in lines[1:]:
        p2 = tf.add_paragraph()
        _set_paragraph(p2, extra, size, color, bold, italic, align, font, spacing)
    return tb


def add_rich(slide, l, t, w, h, runs_spec, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             spacing=None, wrap=True):
    """runs_spec: list of (text, size, color, bold, italic, font) tuples on ONE line."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing is not None:
        p.line_spacing = spacing
    for (text, size, color, bold, italic, font) in runs_spec:
        r = p.add_run()
        r.text = text
        r.font.size = size
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = color
    return tb


def add_bullets(slide, l, t, w, h, items, size=Pt(15), color=TEXT_DARK, font=FONT,
                 bullet_color=ORANGE, gap=Pt(10), bold_lead=None, anchor=MSO_ANCHOR.TOP,
                 marker='—', spacing=1.15):
    """items: list[str] or list[(lead_bold_str, rest_str)]"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = spacing
        p.space_after = gap
        r0 = p.add_run()
        r0.text = f"{marker}  "
        r0.font.size = size
        r0.font.bold = True
        r0.font.name = font
        r0.font.color.rgb = bullet_color
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run()
            r1.text = lead
            r1.font.size = size
            r1.font.bold = True
            r1.font.name = font
            r1.font.color.rgb = color
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = size
            r2.font.bold = False
            r2.font.name = font
            r2.font.color.rgb = color
        else:
            r1 = p.add_run()
            r1.text = item
            r1.font.size = size
            r1.font.bold = False
            r1.font.name = font
            r1.font.color.rgb = color
    return tb


def add_picture(slide, path, l, t, w, h):
    return slide.shapes.add_picture(path, l, t, width=w, height=h)


PAGE_TITLES = {}


def header(slide, kicker, title, page_no, logo_path, title_size=Pt(30)):
    """Consistent top masthead: thin accent rule, small kicker, bold title, logo top-right."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.09), fill=NAVY)
    add_text(slide, MARGIN, Inches(0.34), Inches(9.6), Inches(0.3), kicker.upper(),
              size=Pt(12.5), color=ORANGE, bold=True, font=FONT_HEAD)
    add_text(slide, MARGIN, Inches(0.62), Inches(9.8), Inches(0.7), title,
              size=title_size, color=NAVY_DARK, bold=True, font=FONT_HEAD)
    # logo, top right, fixed height
    lw = Inches(1.55)
    lh = lw * (345 / 1233)
    slide.shapes.add_picture(logo_path, SLIDE_W - MARGIN - lw, Inches(0.42), width=lw, height=lh)
    add_line(slide, MARGIN, Inches(1.28), SLIDE_W - 2 * MARGIN, 0, color=LINE_GRAY, weight=Pt(1))
    PAGE_TITLES[page_no] = title


def footer(slide, page_no, total=12):
    add_text(slide, MARGIN, SLIDE_H - Inches(0.42), Inches(6), Inches(0.3),
              "MoonDAO  |  Deep Space and the Lunar Economy  |  2nd Space Industry Workshop Brazil",
              size=Pt(9), color=TEXT_MUTE, font=FONT)
    add_text(slide, SLIDE_W - MARGIN - Inches(1.2), SLIDE_H - Inches(0.42), Inches(1.2), Inches(0.3),
              f"{page_no:02d} / {total}", size=Pt(9), color=TEXT_MUTE, font=FONT, align=PP_ALIGN.RIGHT)


def stat_card(slide, l, t, w, h, number, label, accent=NAVY):
    add_rect(slide, l, t, w, h, fill=WHITE, line=LINE_GRAY, line_w=Pt(1),
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09, shadow=True)
    add_rect(slide, l, t, Inches(0.07), h, fill=accent, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    add_text(slide, l + Inches(0.22), t + Inches(0.14), w - Inches(0.4), h - Inches(0.7), number,
              size=Pt(28), color=NAVY_DARK, bold=True, font=FONT_HEAD)
    add_text(slide, l + Inches(0.22), t + h - Inches(0.52), w - Inches(0.4), Inches(0.45), label,
              size=Pt(11.5), color=TEXT_GRAY, font=FONT)


def section_number_badge(slide, l, t, text, accent=ORANGE, d=Inches(0.5)):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, d, d)
    shp.fill.solid()
    shp.fill.fore_color.rgb = accent
    _no_line(shp)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.name = FONT_HEAD
    r.font.color.rgb = WHITE
    return shp


def logo_lockup_white(slide, mark_path, l, t, height, mark_ratio=620 / 650):
    """Icon mark + 'MOON'+'DAO' wordmark in white, for dark-background slides."""
    mw = height * mark_ratio
    slide.shapes.add_picture(mark_path, l, t, width=mw, height=height)
    tb = slide.shapes.add_textbox(l + mw + Inches(0.12), t, Inches(3.2), height)
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = "MOON"
    r1.font.size = Pt(26)
    r1.font.bold = True
    r1.font.name = FONT_HEAD
    r1.font.color.rgb = WHITE
    r2 = p.add_run()
    r2.text = "DAO"
    r2.font.size = Pt(13)
    r2.font.bold = True
    r2.font.name = FONT_HEAD
    r2.font.color.rgb = WHITE
    return tb


def add_donut_chart(slide, l, t, w, h, categories, values, colors, hole_size=65):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('Series 1', values)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, l, t, w, h, chart_data)
    chart = gframe.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.has_data_labels = False
    try:
        chart.plots[0].donut_hole_size = hole_size
    except Exception:
        pass
    series = chart.series[0]
    points = series.points
    for i, pt in enumerate(points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colors[i % len(colors)]
        pt.format.line.color.rgb = WHITE
        pt.format.line.width = Pt(2)
    return gframe
